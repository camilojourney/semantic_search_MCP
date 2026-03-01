"""Document text extraction — PDF, DOCX, PPTX.

Converts business documents into structured text pages for the chunking pipeline.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

# Extensions this module handles (binary document formats)
DOCUMENT_EXTENSIONS: set[str] = {".pdf", ".docx", ".pptx"}


@dataclass
class DocumentPage:
    """A page/slide/section of extracted text from a document."""

    text: str
    page_number: int  # 1-indexed
    heading: str | None = None  # section title, slide title, etc.


def is_document(file_path: str | Path) -> bool:
    """Check if a file is a supported document format."""
    return Path(file_path).suffix.lower() in DOCUMENT_EXTENSIONS


def extract_text(file_path: str | Path) -> list[DocumentPage]:
    """Extract text from a document file, routing by extension.

    Returns a list of DocumentPage objects, one per page/slide/section.
    Raises ValueError for unsupported formats.
    """
    file_path = Path(file_path)
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported document format: {ext}")


def _extract_pdf(file_path: Path) -> list[DocumentPage]:
    """Extract text from PDF using pymupdf."""
    import fitz  # pymupdf

    pages: list[DocumentPage] = []
    try:
        doc = fitz.open(str(file_path))
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")
            if text.strip():
                pages.append(DocumentPage(
                    text=text.strip(),
                    page_number=page_num + 1,
                    heading=None,
                ))
        doc.close()
    except Exception as e:
        logger.warning("Failed to extract PDF %s: %s", file_path, e)

    return pages


def _extract_docx(file_path: Path) -> list[DocumentPage]:
    """Extract text from DOCX, grouping by heading sections."""
    from docx import Document

    pages: list[DocumentPage] = []
    try:
        doc = Document(str(file_path))
        current_heading: str | None = None
        current_text: list[str] = []
        section_num = 0

        for para in doc.paragraphs:
            # Detect headings
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                # Save previous section
                if current_text:
                    section_num += 1
                    pages.append(DocumentPage(
                        text="\n".join(current_text).strip(),
                        page_number=section_num,
                        heading=current_heading,
                    ))
                    current_text = []
                current_heading = para.text.strip() or None

            if para.text.strip():
                current_text.append(para.text)

        # Don't forget the last section
        if current_text:
            section_num += 1
            pages.append(DocumentPage(
                text="\n".join(current_text).strip(),
                page_number=section_num,
                heading=current_heading,
            ))

    except Exception as e:
        logger.warning("Failed to extract DOCX %s: %s", file_path, e)

    return pages


def _extract_pptx(file_path: Path) -> list[DocumentPage]:
    """Extract text from PPTX, one page per slide."""
    from pptx import Presentation

    pages: list[DocumentPage] = []
    try:
        prs = Presentation(str(file_path))
        for slide_num, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            slide_title: str | None = None

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

                # Grab slide title
                if hasattr(shape, "is_placeholder") and shape.placeholder_format:
                    if shape.placeholder_format.idx == 0:  # title placeholder
                        slide_title = shape.text.strip() or None

            if texts:
                pages.append(DocumentPage(
                    text="\n".join(texts),
                    page_number=slide_num,
                    heading=slide_title,
                ))

    except Exception as e:
        logger.warning("Failed to extract PPTX %s: %s", file_path, e)

    return pages
